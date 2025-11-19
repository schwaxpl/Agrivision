from pathlib import Path
import os
import sys
try:
    from pptx import Presentation
except Exception as e:
    Presentation = None


PROJECT_ROOT = Path(__file__).resolve().parents[2]
FORMATION_DIR = PROJECT_ROOT / "FormationExistante"
OUTPUT_FILE = FORMATION_DIR / "combined_presentation.md"


def extract_text_from_pptx(pptx_path: Path) -> list:
    """Return a list of slide texts for the given .pptx file."""
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Install with: pip install python-pptx")
    prs = Presentation(str(pptx_path))
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)
        slides_text.append("\n\n".join(parts).strip())
    return slides_text


def gather_presentations(folder: Path) -> dict:
    """Return a dict {filename: [slide_texts,...]} for all .pptx files found."""
    results = {}
    if not folder.exists():
        raise FileNotFoundError(f"Folder not found: {folder}")
    for p in sorted(folder.glob("*.pptx")):
        try:
            slides = extract_text_from_pptx(p)
            results[p.name] = slides
        except Exception as e:
            # skip problematic files but report
            print(f"Warning: failed to read {p.name}: {e}", file=sys.stderr)
    # Optionally warn about .ppt (binary) files which python-pptx doesn't support
    ppt_files = list(folder.glob("*.ppt"))
    if ppt_files:
        names = ", ".join([f.name for f in ppt_files])
        print(f"Note: found .ppt files (not supported): {names}", file=sys.stderr)
    return results


def write_combined_markdown(data: dict, out_file: Path):
    """Write a single markdown file grouping all PPTX contents.

    Structure:
    # Combined Presentations
    ## filename1.pptx
    ### Slide 1
    slide text
    ### Slide 2
    ...
    ## filename2.pptx
    ...
    """
    out_file.parent.mkdir(parents=True, exist_ok=True)
    with out_file.open("w", encoding="utf-8") as f:
        f.write("# Combined Presentations\n\n")
        for fname, slides in data.items():
            f.write(f"## {fname}\n\n")
            if not slides:
                f.write("_No text extracted (possibly slides with only images)_\n\n")
                continue
            for idx, slide in enumerate(slides, start=1):
                f.write(f"### Slide {idx}\n\n")
                # Ensure slide has at least an empty line separation
                f.write(slide + "\n\n")
            f.write("---\n\n")


def main():
    try:
        data = gather_presentations(FORMATION_DIR)
    except FileNotFoundError:
        print(f"Folder not found: {FORMATION_DIR}")
        return
    if not data:
        print("No .pptx files found in FormationExistante.")
        return
    write_combined_markdown(data, OUTPUT_FILE)
    print(f"Combined markdown written to: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
